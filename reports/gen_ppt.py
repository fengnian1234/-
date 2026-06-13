"""
面向民宿客户的精美 PPT 演示文稿
客户 = 民宿老板/运营方，我们是技术搭建方
设计风格：中式简约 · 自然绿意 · 16:9宽屏
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Cm(33.867)   # 16:9
prs.slide_height = Cm(19.05)

# ── 配色 ─────────────────────────────────────────────────
PRIMARY = RGBColor(0x4A, 0x7C, 0x59)
PRIMARY_DARK = RGBColor(0x3B, 0x5E, 0x47)
PRIMARY_LIGHT = RGBColor(0x6B, 0x9B, 0x7E)
ACCENT = RGBColor(0xC8, 0x96, 0x6A)
ACCENT_LIGHT = RGBColor(0xDB, 0xB8, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xFA, 0xFC, 0xF8)
DARK = RGBColor(0x2A, 0x33, 0x28)
GRAY = RGBColor(0x6A, 0x75, 0x68)
LIGHT_GRAY = RGBColor(0x9A, 0xA5, 0x98)
VERY_LIGHT = RGBColor(0xEE, 0xF4, 0xEC)
MIST = RGBColor(0xD4, 0xE5, 0xD8)
DARK_CARD = RGBColor(0x2A, 0x4D, 0x38)
GOLD = RGBColor(0x8B, 0x73, 0x55)
TOTAL = 14

# ── 辅助函数 ─────────────────────────────────────────────

def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=NEAR_WHITE):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def tb(s, l, t, w, h, text, sz=14, c=DARK, bold=False, al=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = bold; p.alignment = al
    return tf

def add_p(tf, text, sz=13, c=DARK, bold=False, al=PP_ALIGN.LEFT, sb=4, sa=4):
    p = tf.add_paragraph(); p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = bold; p.alignment = al
    p.space_before = Pt(sb); p.space_after = Pt(sa)
    return p

def rect(s, l, t, w, h, fill, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line
    else: sh.line.fill.background()
    return sh

def line(s, y, color=MIST):
    rect(s, 3, y, 27.867, 0.03, color)

def grad_bg(s, top_c, bot_c):
    rect(s, 0, 0, 33.867, 9.5, top_c)
    rect(s, 0, 9.5, 33.867, 9.55, bot_c)

def pn(s, n):
    tb(s, 29, 18.0, 4, 0.6, f"{n} / {TOTAL}", 8, LIGHT_GRAY, al=PP_ALIGN.RIGHT)

def section_title(s, icon, title):
    tb(s, 2, 0.8, 25, 1.5, f"{icon}  {title}", 26, PRIMARY_DARK, True)
    line(s, 2.5, ACCENT)

def card(s, x, y, w, h, icon, title, desc, bar_color=PRIMARY):
    rect(s, x, y, w, h, WHITE)
    rect(s, x, y, w, 0.06, bar_color)
    tb(s, x+0.4, y+0.3, w-0.8, 1.2, icon, 22, bar_color, al=PP_ALIGN.CENTER)
    tb(s, x+0.4, y+1.7, w-0.8, 1.0, title, 12, PRIMARY_DARK, True, PP_ALIGN.CENTER)
    tb(s, x+0.4, y+2.9, w-0.8, h-3.1, desc, 9, GRAY, al=PP_ALIGN.CENTER)

def dark_card(s, x, y, w, h, icon, title, items, card_bg=DARK_CARD):
    rect(s, x, y, w, h, card_bg)
    tb(s, x+0.5, y+0.2, w-1, 1.0, icon, 26, WHITE, al=PP_ALIGN.CENTER)
    tb(s, x+0.5, y+1.4, w-1, 0.8, title, 14, WHITE, True, PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        tb(s, x+0.5, y+2.3+j*0.9, w-1, 0.7, f"✓  {item}", 10, RGBColor(0xBB,0xCE,0xC0), al=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
#  S1 — 封面
# ══════════════════════════════════════════════════════════
s = blank_slide(); grad_bg(s, PRIMARY_DARK, PRIMARY)
rect(s, 0, 13, 20, 6.05, RGBColor(0x2A,0x4D,0x38))

tb(s, 2, 3.5, 30, 2.5, "  民宿智能化运营方案", 36, WHITE, True)
tb(s, 2, 6.5, 30, 1.5, "微信公众号客服系统 · AI 智能管家 · 网页全栈搭建", 18, ACCENT_LIGHT)
line(s, 8.5, RGBColor(0x5A,0x8C,0x6A))

tb(s, 2, 9.2, 30, 2.5,
   "为精品民宿提供从「被动等咨询」到「主动服务+自动运营」的智能化方案\n让 AI 管家帮您接待每一位客人，让系统替您完成繁琐的日常运营",
   14, RGBColor(0xCC,0xDC,0xD0))

tb(s, 2, 12.5, 20, 1.5, "  参考案例：庐山 · 云上归墅民宿", 13, ACCENT_LIGHT, True)
tb(s, 2, 14.2, 20, 1.0, "11间客房 | 携程 4.7 分 · 85 条评价 | 系统已上线运行", 11, RGBColor(0xAA,0xBE,0xB0))
tb(s, 2, 16.5, 15, 1.0, "2026年6月", 10, RGBColor(0x99,0xAE,0x9F))
pn(s, 1)


# ══════════════════════════════════════════════════════════
#  S2 — 民宿运营常见困扰
# ══════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_title(s, "", "民宿运营中的常见困扰")

pains = [
    ("*", "咨询量太大\n回复不过来",
     "住客反复问「有什么房型」\n「多少钱」「附近有什么好玩的」\n前台一半时间在做重复回答"),
    ("*", "夜间无人值守\n客人找不到人",
     "前台下班后住客无处求助\n退房后忘记跟进评价\n错过最佳好评时机"),
    ("*", "多渠道订单\n管理混乱",
     "携程、美团、飞猪、大众点评\n每个平台都要登录查看\n哪个渠道卖得好？不知道"),
    ("*", "服务靠记忆\n容易遗漏",
     "谁要打扫？谁要送餐？\n靠记忆、对讲机、纸条\n忙起来总会漏掉一两件"),
]

for i, (icon, title, desc) in enumerate(pains):
    x = 1.2 + i*8.1
    rect(s, x, 4.5, 7.5, 9, WHITE)
    tb(s, x+0.5, 4.8, 6.5, 1.2, icon, 24, PRIMARY, al=PP_ALIGN.CENTER)
    tb(s, x+0.5, 6.2, 6.5, 1.6, title, 12, PRIMARY_DARK, True, PP_ALIGN.CENTER)
    tb(s, x+0.5, 8.2, 6.5, 4.5, desc, 10, GRAY, al=PP_ALIGN.CENTER)

tb(s, 2, 14.5, 30, 1.5, "  我们的方案：用一套系统，解决以上所有问题，让民宿运营从「人工驱动」升级为「系统驱动」", 14, ACCENT, True, PP_ALIGN.CENTER)
pn(s, 2)


# ══════════════════════════════════════════════════════════
#  S3 — 方案总览
# ══════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_title(s, "", "方案总览 — 一套系统，全场景覆盖")

modules = [
    ("[AI]", "AI 智能管家", "三阶段服务\n预订前·预订后·离店后\n7×24h 在线"),
    ("[H5]", "H5 移动端页面", "10 个功能页面\n房型·点单·攻略·服务\n适配所有手机"),
    ("⭐", "积分会员体系", "积分获取/兑换/等级\n激励好评和回头\n增强客户粘性"),
    ("[后台]", "运营管理后台", "多平台订单聚合\n服务请求看板\n经营数据可视化"),
    ("[OTA]", "OTA 平台打通", "携程/美团/飞猪/大众点评\n一键跳转预订\n订单统一管理"),
]
for i, (icon, title, desc) in enumerate(modules):
    card(s, 1.2+i*6.5, 4.2, 6.0, 7.5, icon, title, desc)

line(s, 10, MIST)
tb(s, 2, 10.8, 30, 1.5, "  民宿只需拥有一个微信公众号，其余全部由我们搭建。住客关注即用，无需下载 App。", 13, PRIMARY, True, PP_ALIGN.CENTER)

steps = ["关注公众号", "浏览房型/攻略", "OTA预订", "绑定解锁AI管家", "入住享受服务", "离店好评+积分"]
for i, step in enumerate(steps):
    x = 1.5 + i*5.5
    rect(s, x, 13.5, 5.0, 2.5, VERY_LIGHT)
    tb(s, x+0.3, 13.7, 4.4, 0.8, f"0{i+1}", 18, PRIMARY, True, PP_ALIGN.CENTER)
    tb(s, x+0.3, 14.8, 4.4, 0.8, step, 11, PRIMARY_DARK, al=PP_ALIGN.CENTER)
    if i < 5:
        tb(s, x+5.0, 14.2, 0.5, 0.8, "→", 14, ACCENT, al=PP_ALIGN.CENTER)
pn(s, 3)


# ══════════════════════════════════════════════════════════
#  S4 — AI 管家（重点页）
# ══════════════════════════════════════════════════════════
s = blank_slide(); grad_bg(s, PRIMARY_DARK, PRIMARY)
tb(s, 2, 0.8, 30, 1.2, "  AI 智能管家 — 系统的灵魂", 28, WHITE, True)
tb(s, 2, 2.2, 30, 1.0, "不是关键词回复机器人。是能感知客人阶段、切换服务身份、提供恰当帮助的智能管家。", 13, RGBColor(0xCC,0xDC,0xD0))

phases = [
    ("[1]", "预订前", "免费旅行顾问", "面向所有人\n解答庐山攻略、房型\n自然引导预订", RGBColor(0x2A,0x4D,0x38)),
    ("[2]", "预订后", "专属AI管家", "已预订客人独享\n快捷服务呼叫\n个性化贴心服务", RGBColor(0x3A,0x5D,0x48)),
    ("[3]", "离店后", "复购关怀", "自动好评提醒\n积分+等级通知\n会员回住激励", RGBColor(0x4A,0x6D,0x58)),
]
for i, (icon, phase, title, desc, card_bg) in enumerate(phases):
    x = 1.5 + i*10.8
    rect(s, x, 4.8, 10, 8, card_bg)
    tb(s, x+0.5, 5.2, 9, 1.5, icon, 30, WHITE, al=PP_ALIGN.CENTER)
    tb(s, x+0.5, 6.8, 9, 1.0, phase, 12, ACCENT_LIGHT, True, PP_ALIGN.CENTER)
    tb(s, x+0.5, 7.8, 9, 1.2, title, 18, WHITE, True, PP_ALIGN.CENTER)
    tb(s, x+0.5, 9.2, 9, 3, desc, 12, RGBColor(0xEE,0xF4,0xEC), al=PP_ALIGN.CENTER)

tb(s, 2, 14.5, 30, 1.2, "  一位 AI 管家 = 全年无休的前台 + 旅行顾问 + 售后客服", 15, WHITE, True, PP_ALIGN.CENTER)
tb(s, 2, 16.2, 30, 1.0, "同时服务所有客人，不会疲惫、不会遗漏、不会情绪化。夜间也能即时响应。", 12, RGBColor(0xBB,0xCE,0xC0), al=PP_ALIGN.CENTER)
pn(s, 4)


# ══════════════════════════════════════════════════════════
#  S5 — 房型展示 + 在线点单
# ══════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_title(s, "☕", "房型展示 + 在线点单")

# 房型
rect(s, 1.5, 4.5, 15, 12.5, WHITE)
tb(s, 2, 4.8, 14, 0.8, "  8 种精选房型在线展示", 14, PRIMARY_DARK, True)

rooms = [
    ("特惠单人间 ¥388", "1人 1.5m床 15m²"),
    ("特惠标准间 ¥388", "2人 2×1.2m 15m²"),
    ("知还标准间 ¥488", "2人 2×1.2m 20m²"),
    ("山野大床房 ¥588", "2人 1.8m大床 20m²"),
    ("山景·精致大床房 ¥688", "2人 1.8m 30m² · 全景窗"),
    ("清舍·露台大床房 ¥788", "2人 1.8m 35m² · 独立露台"),
    ("田园家庭房 ¥788", "3人 1.8+1.2m 20m²"),
    ("室雅茶香套房 ¥988", "4人 1.8+1.5m · 独立茶室"),
]
for i, (name, desc) in enumerate(rooms):
    y = 5.9 + i*0.95
    tb(s, 2.0, y, 7, 0.6, name, 9.5, PRIMARY_DARK, True)
    tb(s, 9.0, y, 7, 0.6, desc, 9, GRAY)

# 点单
rect(s, 17.5, 4.5, 15, 12.5, WHITE)
tb(s, 18, 4.8, 14, 0.8, "☕  在线点单 · 饮品定制", 14, PRIMARY_DARK, True)

cats = [
    ("☕ 咖啡+茶饮 (12款)", "美式·拿铁·Dirty·云雾茶…"),
    (" 简餐 (4款)", "意面·披萨·小吃拼盘"),
    (" 甜品 (2款)", "提拉米苏·豆乳蛋糕"),
    (" 特调 (1款)", "「没事的」主理人特调"),
]
for i, (name, items) in enumerate(cats):
    tb(s, 18, 5.9+i*1.1, 14, 0.5, name, 10.5, PRIMARY_DARK, True)
    tb(s, 18, 6.4+i*1.1, 14, 0.5, items, 9, GRAY)

tb(s, 18, 10.5, 14, 0.6, "✨ 饮品定制选项", 12, ACCENT, True)
tb(s, 18, 11.2, 14, 0.6, " 云雾茶：红茶/绿茶 · 杯¥38/壶¥98", 10, DARK)
tb(s, 18, 11.8, 14, 0.6, " 糖度：全糖/半糖/无糖", 10, DARK)
tb(s, 18, 12.4, 14, 0.6, " 冰量：正常冰/少冰/温", 10, DARK)
tb(s, 18, 13.3, 14, 1.5, " 住客房间内浏览菜单\n定制口味 → 微信支付 → 送餐到房\n拓展民宿非房费收入", 10, ACCENT)
pn(s, 5)


# ══════════════════════════════════════════════════════════
#  S6 — 游玩攻略 + 美食推荐
# ══════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_title(s, "", "游玩攻略 + 美食探店")

rect(s, 1.5, 4.2, 15, 13, WHITE)
tb(s, 2, 4.5, 14, 0.8, "  5 条庐山深度游玩路线", 14, PRIMARY_DARK, True)

routes = [
    ("⭐ 一日精华线", "", "6-7h · 西线自然+中线人文"),
    ("⭐ 两日全景（推荐）", "", "2天 · 五老峰+三叠泉"),
    (" 人文寻踪", "", "1天 · 白鹿洞→东林→美庐"),
    (" 好汉坡穿越", "", "5-6h · 15km徒步登顶"),
    (" 半日闲", "", "3-4h · 如琴湖→花径→日落"),
]
for i, (name, diff, desc) in enumerate(routes):
    y = 5.6 + i*1.8
    tb(s, 2, y, 7, 0.6, name, 11, PRIMARY_DARK, True)
    tb(s, 9, y, 1, 0.6, diff, 12, GRAY)
    tb(s, 2, y+0.7, 13, 0.6, desc, 10, GRAY)

rect(s, 17.5, 4.2, 15, 13, WHITE)
tb(s, 18, 4.5, 14, 0.8, "  8 家美食探店推荐", 14, PRIMARY_DARK, True)

foods = [
    " 石牛酒家 · 苍蝇馆子 · ¥55-75",
    " 847别墅餐厅 · 百年别墅 · ¥60-70",
    " 望庐说 · 网红餐厅 · ¥80-100",
    " 庐人村 · 山景餐厅 · ¥80-100",
    " 利民煨汤馆 · 大众9186评 · ¥26-45",
    " 庐山茶饼 · 伴手礼 · ¥10/盒",
    " 庐小仙 · 史努比联名 · ¥18-28",
    " 见山茶 · 网红打卡 · ¥18-32",
]
for i, f in enumerate(foods):
    tb(s, 18, 5.5+i*1.15, 14, 0.7, f, 10, DARK)

tb(s, 18, 15.5, 14, 1.0, " 每条攻略来自真实住客游记\n小红书/大众点评/携程交叉验证\n可跳转高德/腾讯地图导航", 9, LIGHT_GRAY)
pn(s, 6)


# ══════════════════════════════════════════════════════════
#  S7 — 快捷服务 + 积分会员
# ══════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_title(s, "⭐", "快捷服务 + 积分会员")

rect(s, 1.5, 4, 15, 13, WHITE)
tb(s, 2, 4.3, 14, 0.8, "  21 项快捷服务，一键呼叫", 14, PRIMARY_DARK, True)

svcs = [
    (" 客房 (5项)", "续住 · 打扫 · 补充用品 · 送餐 · 洗衣"),
    (" 前台 (13项)", "行李寄存 · 旅游咨询 · 旅拍 · 路线规划\n叫车 · 叫醒 · 搭伙用餐 · 退房\n登山杖 · 雨伞 · 充电宝 · 医药箱 · 特产代购"),
    (" 维修 (3项)", "设施报修 · 空调调节 · 热水问题"),
]
for i, (name, items) in enumerate(svcs):
    y = 5.4 + i*2.8
    tb(s, 2, y, 14, 0.6, name, 11, PRIMARY_DARK, True)
    tb(s, 2, y+0.7, 14, 1.5, items, 9.5, GRAY)

rect(s, 17.5, 4, 15, 13, WHITE)
tb(s, 18, 4.3, 14, 0.8, "⭐ 积分+会员，激励回头客", 14, PRIMARY_DARK, True)

# 积分获取 - 独立文本框
tb(s, 18, 5.2, 14, 0.5, "-- 如何获取积分 --", 11, ACCENT, True)
earn_items = ["入住消费：每消费 1 元 = 1 分",
              "每日签到：+1 分",
              "写 OTA 评价：+80 分（截图发前台）",
              "邀请好友预订：+100 分",
              "生日当月：积分获取 x1.5 倍"]
for i, e in enumerate(earn_items):
    tb(s, 18, 5.8 + i*0.65, 14, 0.55, f"  . {e}", 11, DARK)

# 积分兑换 - 独立文本框
tb(s, 18, 9.3, 14, 0.5, "-- 积分兑换好礼 --", 11, ACCENT, True)
redeem_items = ["精品咖啡 1 杯 · 需 300 分",
                "房型免费升级 · 需 500 分",
                "延迟退房至 14:00 · 需 300 分",
                "房费抵扣券 50 元 · 需 500 分"]
for i, r in enumerate(redeem_items):
    tb(s, 18, 9.9 + i*0.65, 14, 0.55, f"  . {r}", 11, DARK)

# 会员等级 - 独立文本框
tb(s, 18, 12.8, 14, 0.5, "-- 会员等级 --", 11, ACCENT, True)
tier_items = ["银卡(0分起) -> 享 95 折",
              "金卡(累计3000分) -> 享 92 折 + 免费延迟退房",
              "钻石卡(累计8000分) -> 享 9 折 + 免费升级 + 专属管家"]
for i, t in enumerate(tier_items):
    tb(s, 18, 13.4 + i*0.65, 14, 0.55, f"  . {t}", 11, DARK)
pn(s, 7)


# ══════════════════════════════════════════════════════════
#  S8 — 运营后台
# ══════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_title(s, "", "运营后台 — 数据驱动决策")

backends = [
    ("", "多平台订单聚合看板",
     "携程/美团/飞猪/大众点评\n小红书/抖音/直接预订\n7 渠道订单统一管理\n今日入住·退房·在住实时可见\n未来 7 天到店预测",
     "不再需要逐个登录\n各个 OTA 平台查看"),
    ("", "服务请求看板",
     "住客发起的每一条服务请求\n自动记录、状态追踪\n待处理→处理中→已完成\n不再遗漏任何客人需求",
     "不再靠纸条和\n对讲机传话"),
    ("", "经营数据看板",
     "各平台订单数·收入占比\n入住率趋势·月度年度对比\n辅助定价和渠道策略\n每周自动生成经营报告",
     "用数据做决策\n而不是凭感觉"),
    ("⏰", "自动化运营",
     "退房后定时推送好评提醒\n每周自动生成经营周报\n多渠道口碑监控\n系统替你干活，你只看结果",
     "省下时间\n专注于提升服务品质"),
]
for i, (icon, title, desc, value) in enumerate(backends):
    x = 1.2 + i*8.1
    rect(s, x, 4.2, 7.6, 8.5, WHITE)
    tb(s, x+0.4, 4.5, 6.8, 0.8, f"{icon}  {title}", 12, PRIMARY_DARK, True, PP_ALIGN.CENTER)
    tb(s, x+0.4, 5.6, 6.8, 5.0, desc, 10, GRAY, al=PP_ALIGN.CENTER)
    rect(s, x+0.4, 10.8, 6.8, 0.03, ACCENT)
    tb(s, x+0.4, 11.1, 6.8, 1.4, value, 9, ACCENT, True, PP_ALIGN.CENTER)

tb(s, 2, 13.8, 30, 1.2, "  老板在手机上就能查看所有数据，无论身在何处都能掌握民宿运营状况", 13, PRIMARY, True, PP_ALIGN.CENTER)
pn(s, 8)


# ══════════════════════════════════════════════════════════
#  S9 — 系统页面一览
# ══════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_title(s, "", "系统页面一览")

pages = [
    ("", "首页"), ("", "房型列表"), ("", "房型详情"), ("☕", "在线点单"),
    ("", "游玩攻略"), ("", "美食推荐"), ("", "快捷服务"), ("", "地图导航"),
    ("⭐", "积分中心"), ("", "订单管理"), ("", "员工看板"),
]
for i, (icon, name) in enumerate(pages):
    col = i % 4; row = i // 4
    x = 1.5 + col*8.1; y = 4.0 + row*4.0
    rect(s, x, y, 7.5, 3.5, WHITE)
    tb(s, x+0.4, y+0.3, 6.7, 1.2, icon, 22, PRIMARY, al=PP_ALIGN.CENTER)
    tb(s, x+0.4, y+1.8, 6.7, 1.0, name, 12, PRIMARY_DARK, True, PP_ALIGN.CENTER)

tb(s, 2, 16.5, 30, 1.0, "✅  所有页面响应式设计 · 适配各种手机屏幕 · 深色/浅色模式自动切换 · 春夏秋冬三季主题", 10, LIGHT_GRAY, al=PP_ALIGN.CENTER)
pn(s, 9)


# ══════════════════════════════════════════════════════════
#  S10 — 传统 vs 智能化对比
# ══════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_title(s, "", "传统运营 vs 智能化运营")

rect(s, 1.5, 3.8, 31, 1.5, PRIMARY_DARK)
tb(s, 1.5, 4.0, 5, 1.0, "*", 12, WHITE, True, PP_ALIGN.CENTER)
tb(s, 6.5, 4.0, 11, 1.0, "  传统方式", 16, RGBColor(0xFF,0xB0,0xB0), True, PP_ALIGN.CENTER)
tb(s, 17.5, 4.0, 15, 1.0, "✅  智能化方案", 16, RGBColor(0xA0,0xD8,0xA0), True, PP_ALIGN.CENTER)

compares = [
    ("客服响应", "人工回复，晚间无人\n重复回答相同问题", "AI 7×24h 自动应答\n80%+ 咨询无需人工参与"),
    ("房型展示", "住客反复电话询问\n前台口头介绍，信息不一致", "H5 页面自助浏览\n图片+设施+价格清晰透明"),
    ("在线点单", "无线上渠道\n只能到店消费", "手机浏览+口味定制+微信支付\n拓展非房收入"),
    ("好评跟进", "退房后忘记跟进\n人工催评体验尴尬", "自动定时推送+积分激励\n得体高效，转化率更高"),
    ("订单管理", "逐个登录 OTA 平台\n数据分散难以对比", "一站式聚合看板\n各平台数据一目了然"),
    ("会员运营", "无会员体系\n回头客靠缘分", "积分+等级+生日月\n系统化激励复购"),
]

for i, (label, old, new) in enumerate(compares):
    y = 5.6 + i*2.0
    bg_c = VERY_LIGHT if i%2 == 0 else WHITE
    rect(s, 1.5, y, 31, 1.85, bg_c)
    tb(s, 2.0, y+0.3, 4.5, 1.2, label, 12, PRIMARY_DARK, True, PP_ALIGN.CENTER)
    tb(s, 6.5, y+0.3, 11, 1.2, old, 10, GRAY, al=PP_ALIGN.CENTER)
    tb(s, 17.5, y+0.3, 15, 1.2, new, 10, PRIMARY_DARK, al=PP_ALIGN.CENTER)
pn(s, 10)


# ══════════════════════════════════════════════════════════
#  S11 — 案例：云上归墅
# ══════════════════════════════════════════════════════════
s = blank_slide(); grad_bg(s, PRIMARY_DARK, PRIMARY)
tb(s, 2, 0.8, 30, 1.2, "  案例：云上归墅民宿", 28, WHITE, True)
tb(s, 2, 2.2, 30, 1.0, "庐山山上 · 大林沟路27号  |  11间客房  |  携程 4.7 分 · 85 条评价  |  系统已上线", 13, RGBColor(0xCC,0xDC,0xD0))

vals = [
    ("", "AI 管家 7×24h", ["常规咨询自动应答", "前台夜间无需值守", "住客体验一致优质"]),
    ("☕", "在线点单创收", ["咖啡茶饮线上销售", "饮品定制提升体验", "送餐到房无需下楼"]),
    ("⭐", "好评自动跟进", ["退房后智能推送提醒", "积分激励写评价+80", "好评率持续提升"]),
    ("", "会员粘性增强", ["积分+等级+生日月", "钻石卡9折+专属管家", "回头客比例上升"]),
]
for i, (icon, title, items) in enumerate(vals):
    dark_card(s, 1.2+i*8.1, 4.5, 7.6, 7, icon, title, items)

rect(s, 1.2, 12.5, 31.4, 3.5, DARK_CARD)
tb(s, 1.8, 12.8, 30, 1.5, "\"以前住客半夜找前台问第二天去哪玩，我们也答不上来。现在 AI 管家直接给攻略，客人体验好了，我们也省心多了。\"", 13, RGBColor(0xDD,0xE8,0xDD), al=PP_ALIGN.CENTER)
tb(s, 1.8, 14.5, 30, 1.0, "—— 云上归墅民宿运营者", 11, ACCENT_LIGHT, al=PP_ALIGN.RIGHT)
tb(s, 2, 16.8, 30, 1.0, "  云上归墅是我们服务的首个客户，系统已完整上线。欢迎关注公众号「云上归墅民宿」实际体验。", 11, RGBColor(0xAA,0xBE,0xB0), al=PP_ALIGN.CENTER)
pn(s, 11)


# ══════════════════════════════════════════════════════════
#  S12 — 适用民宿类型
# ══════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_title(s, "", "哪些民宿需要这套方案？")

scenarios = [
    ("", "精品民宿/设计师民宿", "5-30 间房，注重品牌和体验\n希望提升调性但无法自建 IT", PRIMARY),
    ("", "景区周边民宿集群", "庐山、大理、莫干山、安吉\n住客对周边攻略需求强烈", ACCENT),
    ("☕", "自带餐饮的民宿", "有咖啡馆/茶室/简餐厅\n希望通过线上拓展非房收入", PRIMARY_LIGHT),
    ("", "多平台运营民宿", "携程/美团/飞猪/大众点评\n同时上线，需统一管理后台", GOLD),
    ("⭐", "希望提升好评的民宿", "入住稳定但评价偏少\n需系统化的评价激励机制", RGBColor(0x5A,0x8C,0x6A)),
]
for i, (icon, title, desc, color) in enumerate(scenarios):
    card(s, 1.2+i*6.6, 4.5, 6.1, 10, icon, title, desc, color)

tb(s, 2, 15.5, 30, 1.5, "  无论您是单体民宿还是小型连锁，只要希望通过技术手段提升运营效率和服务品质，这套方案都适合您。", 13, ACCENT, True, PP_ALIGN.CENTER)
pn(s, 12)


# ══════════════════════════════════════════════════════════
#  S13 — 交付与支持
# ══════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
section_title(s, "", "我们交付什么")

deliverables = [
    ("", "AI 智能管家系统", "三阶段服务：旅行顾问→专属管家→复购关怀\n7×24h 自动值守，住客关注公众号即可使用"),
    ("", "移动端 H5 网页", "10 个功能页面，嵌入公众号菜单\n房型·点单·攻略·服务·积分中心\n响应式设计，适配所有手机屏幕"),
    ("⭐", "积分会员体系", "积分获取/兑换/等级管理\n激励好评 + 回头入住\n生日月 ×1.5 倍积分加成"),
    ("[后台]", "运营管理后台", "多平台订单聚合看板\n服务请求追踪看板\n经营数据 + 每周自动周报"),
    ("", "预订渠道对接", "携程/美团/飞猪/大众点评\n一键跳转预订 + 订单统一管理\n退房后好评自动推送"),
    ("", "视觉定制", "匹配民宿品牌配色和设计风格\n深色/浅色模式 + 季节主题\n民宿 Logo 和视觉元素融入"),
]

for i, (icon, title, desc) in enumerate(deliverables):
    col = i % 3; row = i // 3
    x = 1.2 + col*10.8; y = 4.2 + row*6.2
    rect(s, x, y, 10, 5.5, WHITE)
    tb(s, x+0.5, y+0.2, 9, 1.0, icon, 24, PRIMARY, al=PP_ALIGN.CENTER)
    tb(s, x+0.5, y+1.4, 9, 0.8, title, 13, PRIMARY_DARK, True, PP_ALIGN.CENTER)
    tb(s, x+0.5, y+2.4, 9, 2.5, desc, 10, GRAY, al=PP_ALIGN.CENTER)

pn(s, 13)


# ══════════════════════════════════════════════════════════
#  S14 — 尾页
# ══════════════════════════════════════════════════════════
s = blank_slide(); grad_bg(s, PRIMARY_DARK, PRIMARY)

tb(s, 2, 3.5, 30, 2, "  让您的民宿也拥有智能大脑", 34, WHITE, True, PP_ALIGN.CENTER)
tb(s, 2, 6.5, 30, 1.5, "微信公众号客服系统 · AI 智能管家 · 在线点单 · 会员积分 · 订单管理", 16, ACCENT_LIGHT, al=PP_ALIGN.CENTER)
line(s, 8.5, RGBColor(0x5A,0x8C,0x6A))

tb(s, 2, 9.5, 30, 2, "我们提供的不是一个模板，而是一套为民宿量身定制的完整解决方案\n以云上归墅为真实案例，已在庐山风景区实际运行", 13, RGBColor(0xCC,0xDC,0xD0), al=PP_ALIGN.CENTER)

rect(s, 8, 12.5, 18, 3.5, DARK_CARD)
tb(s, 8.5, 12.8, 17, 1.5, "  欢迎关注「云上归墅民宿」微信公众号\n实际体验全套系统功能", 13, WHITE, True, PP_ALIGN.CENTER)
tb(s, 8.5, 14.5, 17, 1.0, "微信中搜索「云上归墅民宿」", 11, RGBColor(0xBB,0xCE,0xC0), al=PP_ALIGN.CENTER)

tb(s, 2, 17, 30, 1.0, "方案版本 v2.14  ·  2026年6月", 10, RGBColor(0x99,0xAE,0x9F), al=PP_ALIGN.CENTER)
pn(s, 14)


# ── 保存 ─────────────────────────────────────────────────
output = os.path.join(os.path.dirname(__file__), "云上归墅_民宿智能化运营方案_v2_20260613.pptx")
prs.save(output)
print(f"✅ PPT 已保存至: {output}")
print(f" 共 {len(prs.slides)} 页")
